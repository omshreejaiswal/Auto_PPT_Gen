from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import math


def parse_outline(slide_text):
    slides = []
    blocks = slide_text.split("Slide Title:")

    for block in blocks:
        if not block.strip():
            continue

        lines = [l.strip() for l in block.split("\n") if l.strip()]
        title = lines[0]
        bullets = [l.replace("-", "").strip() for l in lines[1:]]

        slides.append({
            "title": title,
            "bullets": bullets
        })

    return slides


def estimate_slide_count(text, bullets_per_slide=3):
    paragraphs = [p for p in text.split("\n") if p.strip()]
    return max(1, math.ceil(len(paragraphs) / bullets_per_slide))


def find_best_content_layout(prs):
    """
    Pick a template layout that already has title + body.
    """
    for layout in prs.slide_layouts:
        has_title = False
        has_body = False

        for shape in layout.shapes:
            if shape.is_placeholder:
                if shape.placeholder_format.type == 1:
                    has_title = True
                if shape.placeholder_format.type == 2:
                    has_body = True

        if has_title and has_body:
            return layout

    return prs.slide_layouts[1]  # safe fallback


def improve_existing_slides(template_file, slide_text, raw_text):
    prs = Presentation(template_file)

    structured_slides = parse_outline(slide_text)
    content_layout = find_best_content_layout(prs)

    existing_slides = list(prs.slides)
    slide_ptr = 0

    for data in structured_slides:
        # 1️⃣ Improve existing slide
        if slide_ptr < len(existing_slides):
            slide = existing_slides[slide_ptr]
        # 2️⃣ Append new slide
        else:
            slide = prs.slides.add_slide(content_layout)

        # ---- TITLE ----
        if slide.shapes.title:
            slide.shapes.title.text = data["title"]

        # ---- BODY ----
        body_found = False
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                tf = shape.text_frame
                tf.clear()
                for bullet in data["bullets"]:
                    tf.add_paragraph().text = bullet
                body_found = True
                break

        # If no body placeholder exists → add textbox
        if not body_found:
            textbox = slide.shapes.add_textbox(
                prs.slide_width * 0.1,
                prs.slide_height * 0.25,
                prs.slide_width * 0.8,
                prs.slide_height * 0.6,
            )
            tf = textbox.text_frame
            tf.clear()
            for bullet in data["bullets"]:
                tf.add_paragraph().text = bullet

        slide_ptr += 1

    return prs


# Backward compatibility
create_ppt = improve_existing_slides